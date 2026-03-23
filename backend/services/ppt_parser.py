from pptx import Presentation

def extract_text_from_pptx(file_path: str):
    """
    Extracts text from a PPT/PPTX file slide by slide.
    """
    slides_text = []
    
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            
            # Optionally extract speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    slide_text += "\n[Notes:] " + notes_slide.notes_text_frame.text
                        
            cleaned_text = " ".join(slide_text.split()).strip()
            if cleaned_text:
                slides_text.append(cleaned_text)
                
    except Exception as e:
        print(f"Error parsing PPTX: {e}")
        
    return slides_text
